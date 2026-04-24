#!/usr/bin/env python3
"""
doc_parser.py: 文档解析模块
支持 Excel/Word/PDF/PPT 格式解析
"""

import os
import json
import tempfile
import subprocess
from pathlib import Path
from typing import List, Dict, Tuple, Optional

# 尝试导入可选库
try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def get_file_token_from_doc(doc_token: str) -> Tuple[str, str]:
    """根据 doc_token 获取 file_token

    返回: (file_token, error_msg)
    """
    cmd = f'lark-cli wiki spaces get_node --params \'{{"token":"{doc_token}"}}\''
    result = subprocess.run(
        cmd,
        shell=True,
        capture_output=True,
        text=True
    )
    if result.returncode != 0:
        return "", f"获取文件元数据失败: {result.stderr}"

    try:
        data = json.loads(result.stdout)
        if data.get("code") == 0:
            node = data.get("data", {}).get("node", {})
            file_token = node.get("obj_token", "") or node.get("file_token", "")
            if file_token:
                return file_token, ""
            return "", "未找到 obj_token 或 file_token"
        return "", f"API 错误: {data}"
    except json.JSONDecodeError:
        return "", f"解析元数据失败: {result.stdout}"


def download_file(file_token: str, output_filename: str) -> Tuple[bool, str]:
    """使用 lark-cli 下载文件到本地

    Args:
        file_token: 文件 token
        output_filename: 相对文件名（不含路径）

    Returns: (success, error_msg)
    """
    cmd = f'lark-cli drive +download --file-token "{file_token}" --output "{output_filename}"'
    result = subprocess.run(
        cmd,
        shell=True,
        capture_output=True,
        text=True
    )
    return result.returncode == 0, result.stderr if result.returncode != 0 else ""


def cleanup_file(filename: str):
    """清理下载的文件"""
    try:
        if os.path.exists(filename):
            os.remove(filename)
    except:
        pass


def parse_excel(file_path: str, max_chars: int = 50000) -> Tuple[str, Optional[str]]:
    """解析 Excel 文件

    Returns: (content, error_msg)
    """
    if not HAS_OPENPYXL:
        return "", "未安装 openpyxl，请运行: pip install openpyxl"

    try:
        text_parts = []
        wb = openpyxl.load_workbook(file_path, data_only=True)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"\n\n=== 工作表: {sheet_name} ===\n")

            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)

                if sum(len(t) for t in text_parts) > max_chars:
                    break

            if sum(len(t) for t in text_parts) > max_chars:
                break

        content = "\n".join(text_parts)[:max_chars]
        return content, None
    except Exception as e:
        return "", f"Excel 解析失败: {e}"


def parse_word(file_path: str, max_chars: int = 50000) -> Tuple[str, Optional[str]]:
    """解析 Word 文件

    Returns: (content, error_msg)
    """
    if not HAS_DOCX:
        return "", "未安装 python-docx，请运行: pip install python-docx"

    try:
        doc = Document(file_path)
        paragraphs = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
                if sum(len(p) for p in paragraphs) > max_chars:
                    break

        content = "\n\n".join(paragraphs)[:max_chars]
        return content, None
    except Exception as e:
        return "", f"Word 解析失败: {e}"


def parse_pdf(file_path: str, max_chars: int = 50000) -> Tuple[str, Optional[str]]:
    """解析 PDF 文件

    Returns: (content, error_msg)
    """
    if not HAS_PDFPLUMBER:
        return "", "未安装 pdfplumber，请运行: pip install pdfplumber"

    try:
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"\n\n=== 第 {page_num + 1} 页 ===\n")
                    text_parts.append(page_text)
                    if sum(len(t) for t in text_parts) > max_chars:
                        break

        content = "".join(text_parts)[:max_chars]
        return content, None
    except Exception as e:
        return "", f"PDF 解析失败: {e}"


def parse_pptx(file_path: str, max_chars: int = 50000) -> Tuple[str, Optional[str]]:
    """解析 PowerPoint 文件

    Returns: (content, error_msg)
    """
    if not HAS_PPTX:
        return "", "未安装 python-pptx，请运行: pip install python-pptx"

    try:
        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides):
            text_parts.append(f"\n\n=== 幻灯片 {slide_num + 1} ===\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
                    text_parts.append("\n")

            if sum(len(t) for t in text_parts) > max_chars:
                break

        content = "".join(text_parts)[:max_chars]
        return content, None
    except Exception as e:
        return "", f"PPTX 解析失败: {e}"


def parse_csv(file_path: str, max_chars: int = 50000) -> Tuple[str, Optional[str]]:
    """解析 CSV 文件

    Returns: (content, error_msg)
    """
    try:
        import pandas as pd
        df = pd.read_csv(file_path)
        content = df.to_string(max_chars=max_chars)
        return content[:max_chars], None
    except Exception as e:
        return "", f"CSV 解析失败: {e}"


def parse_document_from_token(doc_token: str, doc_type: str, max_chars: int = 50000) -> Tuple[str, Optional[str]]:
    """从飞书 token 解析文档

    Args:
        doc_token: 文档 token
        doc_type: 文档类型 (xlsx/xls/docx/pdf/csv/pptx)

    Returns: (content, error_msg)
    """
    # 先获取 file_token
    file_token, err = get_file_token_from_doc(doc_token)
    if err:
        return "", f"获取 file_token 失败: {err}"

    # 根据类型解析
    file_ext = doc_type.lower()
    if file_ext in ("xlsx", "xls"):
        output_filename = f"{file_token}.xlsx"
        parse_func = parse_excel
    elif file_ext == "pdf":
        output_filename = f"{file_token}.pdf"
        parse_func = parse_pdf
    elif file_ext in ("docx", "doc"):
        output_filename = f"{file_token}.docx"
        parse_func = parse_word
    elif file_ext == "pptx":
        output_filename = f"{file_token}.pptx"
        parse_func = parse_pptx
    elif file_ext == "csv":
        output_filename = f"{file_token}.csv"
        parse_func = parse_csv
    else:
        return "", f"不支持的文档类型: {doc_type}"

    success, err = download_file(file_token, output_filename)
    if not success:
        return "", f"下载失败: {err}"

    try:
        content, err = parse_func(output_filename, max_chars)
        return content, err
    finally:
        cleanup_file(output_filename)


def parse_native_docx(content: str, max_chars: int = 50000) -> str:
    """解析飞书原生文档（Markdown 格式）

    直接返回 Markdown 内容，无需额外解析
    """
    return content[:max_chars]


def get_supported_types() -> Dict[str, str]:
    """获取支持的文档类型"""
    return {
        "xlsx": "Excel",
        "xls": "Excel",
        "docx": "Word",
        "doc": "Word",
        "pdf": "PDF",
        "pptx": "PowerPoint",
        "csv": "CSV"
    }


def main():
    """测试解析功能"""
    print("=== doc_parser.py 测试 ===\n")
    print("支持的文档类型:")
    for ext, name in get_supported_types().items():
        print(f"  .{ext}: {name}")
    print()
    print("库状态:")
    print(f"  openpyxl: {'已安装' if HAS_OPENPYXL else '未安装'}")
    print(f"  python-docx: {'已安装' if HAS_DOCX else '未安装'}")
    print(f"  pdfplumber: {'已安装' if HAS_PDFPLUMBER else '未安装'}")
    print(f"  python-pptx: {'已安装' if HAS_PPTX else '未安装'}")


if __name__ == "__main__":
    main()